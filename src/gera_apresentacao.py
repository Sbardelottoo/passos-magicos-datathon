"""Gera a apresentacao gerencial (storytelling) em PPTX.

Padrao de apresentacao executiva:
  - Titulos ASSERTIVOS (a conclusao no titulo, nao um rotulo)
  - Resumo executivo logo no inicio
  - Numeros grandes (KPI hero) em vez de numeros dentro de bullets
  - Identidade visual da Passos Magicos (azul/dourado + logo)
  - Bullets curtos (1 linha sempre que possivel)
"""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).resolve().parents[1]
FIG = BASE / "reports" / "figures"
# PowerPoint/PDF renderizam PNG com alpha de forma inconsistente (viram fundo
# preto). Por isso usamos duas variantes ja achatadas sobre cor solida:
LOGO = BASE / "streamlit_app" / "assets" / "logo_branco.png"    # p/ fundo branco
LOGO_NAVY = BASE / "streamlit_app" / "assets" / "logo_navy.png"  # p/ fundo azul-escuro
OUT = BASE / "apresentacao" / "storytelling_passos_magicos.pptx"
OUT.parent.mkdir(exist_ok=True)

APP_URL = "paapps-magicos-datathon-mbjdeucxaqr7pkuodmftr6.streamlit.app"

# ---- Paleta (alinhada ao app e ao logo) ----
AZUL_ESCURO = RGBColor(0x1F, 0x3A, 0x5F)   # faixa de titulo / fundos
AZUL = RGBColor(0x0A, 0x5F, 0xA8)          # primario
DOURADO = RGBColor(0xE9, 0xA2, 0x3B)       # acento / destaque
CINZA = RGBColor(0x33, 0x3F, 0x4C)         # texto corrente
CINZA_CLARO = RGBColor(0x8A, 0x94, 0xA0)   # texto secundario
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
VERDE = RGBColor(0x2E, 0x9E, 0x5B)
VERMELHO = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_n = {"i": 0}  # contador de slides


# --------------------------------------------------------------- helpers
def cx_texto(slide, x, y, w, h, texto, tam=18, cor=CINZA, bold=False,
             align=PP_ALIGN.LEFT, italico=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(tam)
    r.font.bold = bold
    r.font.italic = italico
    r.font.color.rgb = cor
    return tb


def bullets(slide, x, y, w, h, itens, tam=16):
    """itens: lista de (texto, destaque_bool)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, destaque) in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = "▪  " + txt
        r.font.size = Pt(tam)
        r.font.bold = destaque
        # destaque em AZUL (nao mais roxo); texto normal em cinza
        r.font.color.rgb = AZUL if destaque else CINZA
    return tb


def rodape(slide, com_logo=True):
    """Numero do slide + logo discreto no rodape."""
    _n["i"] += 1
    cx_texto(slide, 0.5, 6.95, 1.0, 0.35, str(_n["i"]), 11, CINZA_CLARO)
    if com_logo and LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(12.2), Inches(6.75), height=Inches(0.55))


def faixa_titulo(slide, titulo, subtitulo=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15))
    s.fill.solid()
    s.fill.fore_color.rgb = AZUL_ESCURO
    s.line.fill.background()
    cx_texto(slide, 0.5, 0.18, 12.3, 0.6, titulo, 25, BRANCO, True)
    if subtitulo:
        cx_texto(slide, 0.5, 0.72, 12.3, 0.35, subtitulo, 13, RGBColor(0xC5, 0xD3, 0xE8))


def add_picture_fit(slide, caminho, x, y, max_w, max_h):
    """Insere a imagem respeitando a proporcao real, sem nunca estourar a
    caixa (max_w x max_h) - evita cortar a figura na borda do slide."""
    with Image.open(caminho) as im:
        aspecto = im.width / im.height
    if max_w / max_h > aspecto:
        h, w = max_h, max_h * aspecto
    else:
        w, h = max_w, max_w / aspecto
    # centraliza na caixa disponivel
    cx = x + (max_w - w) / 2
    cy = y + (max_h - h) / 2
    slide.shapes.add_picture(str(caminho), Inches(cx), Inches(cy), Inches(w), Inches(h))


def slide_conteudo(titulo, itens, imagem=None, subtitulo=None, tam_bullet=16):
    s = prs.slides.add_slide(BLANK)
    faixa_titulo(s, titulo, subtitulo)
    if imagem and (FIG / imagem).exists():
        bullets(s, 0.6, 1.55, 5.3, 5.0, itens, tam_bullet)
        add_picture_fit(s, FIG / imagem, 6.1, 1.5, 6.75, 4.95)
    else:
        bullets(s, 0.7, 1.6, 12.0, 5.0, itens, tam_bullet)
    rodape(s)
    return s


def slide_secao(titulo, subtitulo=""):
    s = prs.slides.add_slide(BLANK)
    fundo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = AZUL_ESCURO
    fundo.line.fill.background()
    cx_texto(s, 0.9, 2.9, 11.5, 1.2, titulo, 38, BRANCO, True)
    if subtitulo:
        cx_texto(s, 0.9, 4.1, 11.5, 0.8, subtitulo, 19, DOURADO)
    _n["i"] += 1
    return s


def kpi(slide, x, valor, rotulo, detalhe, cor=AZUL, de=None):
    """KPI em bloco vertical, sem risco de sobreposicao:
    [de: opcional, pequeno] -> [valor: grande, 1 linha] -> rotulo -> detalhe.
    """
    y = 2.15
    if de:
        cx_texto(slide, x, y, 3.9, 0.4, de, 16, CINZA_CLARO, False, PP_ALIGN.CENTER)
    y += 0.5  # espaco reservado sempre, com ou sem "de", p/ alinhar as colunas
    cx_texto(slide, x, y, 3.9, 0.95, valor, 48, cor, True, PP_ALIGN.CENTER)
    y += 1.05
    cx_texto(slide, x, y, 3.9, 0.4, rotulo, 15, CINZA, True, PP_ALIGN.CENTER)
    cx_texto(slide, x, y + 0.42, 3.9, 0.8, detalhe, 12, CINZA_CLARO, False, PP_ALIGN.CENTER)


# =========================================================== 1. CAPA
s = prs.slides.add_slide(BLANK)
fundo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
fundo.fill.solid()
fundo.fill.fore_color.rgb = AZUL_ESCURO
fundo.line.fill.background()
if LOGO_NAVY.exists():
    s.shapes.add_picture(str(LOGO_NAVY), Inches(0.9), Inches(1.25), height=Inches(1.5))
cx_texto(s, 0.9, 3.05, 11.5, 1.2, "Educação que transforma", 44, BRANCO, True)
cx_texto(s, 0.9, 4.15, 11.5, 1.0,
         "Como os dados da PEDE mostram o impacto do programa — e antecipam\n"
         "quais alunos precisam de ajuda antes da queda", 19, RGBColor(0xC5, 0xD3, 0xE8))
faixa = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.75), Inches(3.2), Inches(0.06))
faixa.fill.solid()
faixa.fill.fore_color.rgb = DOURADO
faixa.line.fill.background()
cx_texto(s, 0.9, 6.1, 11.5, 0.8,
         "Associação Passos Mágicos · PEDE 2022–2024\nDatathon PosTech — Fase 5",
         13, CINZA_CLARO)

# =========================================================== 2. MENSAGENS-CHAVE
s = prs.slides.add_slide(BLANK)
faixa_titulo(s, "Três mensagens, se você só ler este slide",
             "Resumo executivo")
msgs = [
    ("1", "O programa funciona.",
     "A defasagem moderada/severa caiu de 22,2% para 8,0% em dois anos, e a melhora "
     "se confirma na coorte fixa — não é troca de alunos."),
    ("2", "Engajamento é a maior alavanca.",
     "IEG e IDA puxam juntos o ponto de virada e a nota global (INDE). O efeito é "
     "cumulativo: agir em vários eixos rende muito mais."),
    ("3", "Agora dá para agir antes da queda.",
     "Um modelo preditivo identifica ~3 em cada 4 alunos que piorariam no ano seguinte, "
     "entregue como aplicação pronta para a equipe usar."),
]
y = 1.6
for num, titulo, texto in msgs:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(y), Inches(0.62), Inches(0.62))
    circ.fill.solid()
    circ.fill.fore_color.rgb = DOURADO
    circ.line.fill.background()
    cx_texto(s, 0.75, y + 0.09, 0.62, 0.45, num, 20, BRANCO, True, PP_ALIGN.CENTER)
    cx_texto(s, 1.65, y - 0.02, 10.9, 0.45, titulo, 20, AZUL, True)
    cx_texto(s, 1.65, y + 0.48, 10.9, 0.9, texto, 14, CINZA)
    y += 1.72
rodape(s)

# =========================================================== 3. CONTEXTO
slide_conteudo(
    "O desafio: transformar 3 anos de avaliações em decisão pedagógica",
    [("A Passos Mágicos usa a educação para mudar a vida de crianças em vulnerabilidade.", False),
     ("A PEDE mede 7 indicadores (IAN, IDA, IEG, IAA, IPS, IPP, IPV), sintetizados no INDE.", False),
     ("Base analisada: 3.030 registros aluno-ano · 1.661 alunos · 2022 a 2024.", True),
     ("A chave RA permite acompanhar o mesmo aluno ao longo dos anos.", True),
     ("Nota metodológica: usamos a base real (abas 2022/2023/2024); o dicionário fornecido "
      "descreve um formato antigo e serviu como referência conceitual.", False)],
    subtitulo="Contexto e dados")

# =========================================================== 4. KPI HERO
s = prs.slides.add_slide(BLANK)
faixa_titulo(s, "Dois anos de programa, em três números", "O impacto medido pela PEDE")
kpi(s, 0.7, "8,0%", "Defasagem moderada/severa", "queda de 14 pontos percentuais", VERDE, de="2022: 22,2%  →  2024:")
kpi(s, 4.72, "68,0%", "Alunos em pedras superiores", "Ametista + Topázio", AZUL, de="2022: 55,6%  →  2024:")
kpi(s, 8.74, "29%", "Subiram de pedra", "na coorte fixa de 468 alunos (2022→2024)", DOURADO)
linha = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.03))
linha.fill.solid()
linha.fill.fore_color.rgb = RGBColor(0xE0, 0xE5, 0xEA)
linha.line.fill.background()
cx_texto(s, 0.9, 5.6, 11.5, 0.9,
         "A coorte fixa (mesmos alunos nos três anos) confirma que a melhora é real — "
         "não é efeito da entrada e saída de alunos.", 15, CINZA, italico=True)
rodape(s)

# =========================================================== 5. IAN
slide_conteudo(
    "A defasagem caiu pela metade — e mais que dobrou o grupo adequado",
    [("Defasagem moderada/severa: 22,2% (2022) → 8,0% (2024).", True),
     ("Alunos adequados ou adiantados: 30,1% → 53,8%.", True),
     ("Ainda restam 93 alunos em defasagem moderada/severa.", False),
     ("→ Esse é exatamente o grupo que o modelo preditivo ajuda a priorizar.", False)],
    "q1_defasagem_por_ano.png", subtitulo="Adequação ao nível (IAN)")

# =========================================================== 6. PEDRAS
slide_conteudo(
    "Os alunos sobem de nível: 130 mudaram de pedra para melhor",
    [("Pedras superiores (Ametista + Topázio): 55,6% → 68,0%.", True),
     ("Na coorte fixa, 130 alunos subiram de pedra entre 2022 e 2024.", True),
     ("A matriz de transição mostra mobilidade ascendente consistente.", False),
     ("→ Evidência direta do impacto do programa, aluno a aluno.", False)],
    "q10_transicao_pedras.png", subtitulo="Efetividade por pedra (Quartzo → Topázio)")

# =========================================================== 7. IDA / INDE
slide_conteudo(
    "Desempenho e nota global crescem de forma consistente",
    [("INDE médio: 7,04 → 7,34 → 7,40.", True),
     ("IDA médio: 6,09 (2022) → 6,35 (2024).", True),
     ("Fases iniciais têm IDA mais alto; há queda nas fases de transição.", False),
     ("→ As transições de fase merecem acompanhamento reforçado.", False)],
    "q10_evolucao_indicadores.png", subtitulo="Desempenho acadêmico (IDA) e INDE")

# =========================================================== 8. DRIVERS
slide_conteudo(
    "Engajar o aluno move desempenho, ponto de virada e nota global juntos",
    [("IEG × IDA: r = 0,54 · IEG × IPV: r = 0,56.", True),
     ("Maiores alavancas do INDE: IDA (0,79), IEG (0,75), IPV (0,72).", True),
     ("Principais motores do IPV: IPP (0,61), IEG (0,56), IDA (0,56).", False),
     ("→ Engajamento é o ponto de entrada de maior retorno.", False)],
    "q8_drivers_inde.png", subtitulo="O que move os resultados")

# =========================================================== 9. COMBINAÇÕES
slide_conteudo(
    "O ganho é cumulativo: agir em vários eixos vale mais que em um só",
    [("Com IDA+IEG+IPS+IPP acima da mediana, o INDE médio vai a 8,42.", True),
     ("Quando nenhum está alto, o INDE médio fica em 6,04.", True),
     ("A regressão recupera os pesos da fórmula do INDE (R²=1,0): IDA e IEG pesam mais.", False),
     ("→ Intervenções combinadas rendem mais que esforço isolado.", False)],
    "q8_combinacoes_inde.png", subtitulo="Multidimensionalidade dos indicadores")

# =========================================================== 10. ALERTAS
slide_conteudo(
    "O psicossocial cai antes do desempenho — é o alerta mais precoce",
    [("Quem piora no ano seguinte já tinha IPS mais baixo: 5,80 vs 6,39.", True),
     ("Autoavaliação descolada da realidade: 444 alunos com IAA≥8 e IDA≤5.", True),
     ("IPP capta dimensão complementar ao IAN (r≈0,12), não redundante.", False),
     ("→ Monitorar IPS funciona como termômetro antecipado.", False)],
    "q5_ips_antecede_queda.png", subtitulo="Sinais de alerta (IPS, IAA, IPP)")

# =========================================================== 11. SEÇÃO MODELO
slide_secao("Do diagnóstico à antecipação",
            "Um modelo que aponta quem vai piorar — antes de piorar")

# =========================================================== 12. COMO FUNCIONA
slide_conteudo(
    "Prevemos a mudança futura, não o retrato de hoje",
    [("Alvo: probabilidade de a defasagem PIORAR no ano seguinte.", True),
     ("Usa os indicadores atuais do aluno + idade, fase e tempo de casa.", False),
     ("Feature engineering longitudinal via RA · split estratificado 75/25.", False),
     ("Três algoritmos comparados; vencedor: Random Forest.", False),
     ("Cuidado com vazamento: prever a defasagem atual seria trivial "
      "(ela deriva de idade e fase) — por isso prevemos a mudança.", True)],
    subtitulo="Como o modelo foi construído")

# =========================================================== 13. RESULTADOS
s = prs.slides.add_slide(BLANK)
faixa_titulo(s, "O modelo antecipa 3 em cada 4 alunos que vão piorar",
             "Resultados e transparência")
kpi(s, 0.7, "0,88", "AUC (teste)", "capacidade de separar quem piora", AZUL)
kpi(s, 4.72, "73%", "Recall", "dos que pioram, quantos capturamos", VERDE)
kpi(s, 8.74, "51%", "Precisão", "dos sinalizados, quantos de fato pioram", DOURADO)
cx_texto(s, 0.9, 5.35, 11.5, 1.2,
         "Assumindo a limitação: cerca de metade dos alunos sinalizados são falsos alarmes. "
         "Por isso a ferramenta é de TRIAGEM — ela ordena a fila de atenção, e a decisão "
         "final continua sendo da equipe pedagógica.", 15, CINZA, italico=True)
rodape(s)

# =========================================================== 14. PREDITORES
slide_conteudo(
    "Entre os indicadores pedagógicos, o ponto de virada é o que mais antecipa",
    [("Peso maior: defasagem atual, IAN, idade e fase (posição no ciclo).", False),
     ("Entre os pedagógicos, IPV é o mais preditivo — seguido de IDA.", True),
     ("Leitura de negócio: manter o aluno rumo ao ponto de virada protege.", True),
     ("→ Reforço de aprendizagem deve andar junto com suporte socioemocional.", False)],
    "modelo_importancias.png", subtitulo="O que mais pesa na previsão")

# =========================================================== 15. A FERRAMENTA
s = prs.slides.add_slide(BLANK)
faixa_titulo(s, "Entregue como ferramenta, não como relatório",
             "Aplicação publicada e pronta para uso")
bullets(s, 0.7, 1.6, 7.4, 4.6, [
    ("Simulador: risco de um perfil, com o PORQUÊ e as ações sugeridas.", True),
    ("Turma priorizada: envie a planilha PEDE e receba a fila de atenção.", True),
    ("Histórico do aluno: trajetória 2022–2024 e risco atual.", False),
    ("Panorama: indicadores, pedras e evidências de impacto.", False),
    ("Sobre o modelo: desempenho, limitações e uso responsável.", False),
    ("Relatório em PDF por aluno, pronto para a reunião pedagógica.", False),
], 15)
cx_texto(s, 8.4, 2.3, 4.4, 0.4, "Acesse a aplicação", 15, CINZA, True, PP_ALIGN.CENTER)
cx_texto(s, 8.3, 2.85, 4.6, 1.2, APP_URL, 12, AZUL, True, PP_ALIGN.CENTER)
cx_texto(s, 8.4, 4.15, 4.4, 0.8,
         "Publicada no Streamlit\nCommunity Cloud", 12, CINZA_CLARO, False, PP_ALIGN.CENTER)
rodape(s)

# =========================================================== 16. RECOMENDAÇÕES
s = prs.slides.add_slide(BLANK)
faixa_titulo(s, "Três prioridades para a Passos Mágicos", "Recomendações")
recs = [
    ("Triagem anual com o modelo",
     "Rodar a lista priorizada no início do ciclo e direcionar o acompanhamento "
     "para quem tem maior risco."),
    ("IPS como termômetro precoce",
     "Acompanhar quedas psicossociais: elas antecedem as quedas acadêmicas em um ano."),
    ("Investir em engajamento (IEG/IPV)",
     "É a alavanca de maior retorno sobre a nota global — e protege contra a defasagem."),
]
y = 1.65
for i, (titulo, texto) in enumerate(recs, 1):
    barra = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(y), Inches(0.09), Inches(0.95))
    barra.fill.solid()
    barra.fill.fore_color.rgb = DOURADO
    barra.line.fill.background()
    cx_texto(s, 1.1, y - 0.05, 11.4, 0.45, f"{i}. {titulo}", 20, AZUL, True)
    cx_texto(s, 1.1, y + 0.42, 11.4, 0.6, texto, 14, CINZA)
    y += 1.35
cx_texto(s, 0.75, 5.85, 11.9, 1.0,
         "Também recomendamos: atenção às fases de transição (onde o IDA cai), trabalho de "
         "autopercepção com alunos de IAA alto e IDA baixo, e padronização da coleta "
         "(fase, gênero, IPP) para fortalecer as próximas análises.", 13, CINZA_CLARO)
rodape(s)

# =========================================================== 17. FECHAMENTO
s = prs.slides.add_slide(BLANK)
fundo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
fundo.fill.solid()
fundo.fill.fore_color.rgb = AZUL_ESCURO
fundo.line.fill.background()
if LOGO_NAVY.exists():
    s.shapes.add_picture(str(LOGO_NAVY), Inches(0.9), Inches(1.1), height=Inches(1.2))
cx_texto(s, 0.9, 2.6, 11.5, 0.9, "Próximos passos", 36, BRANCO, True)
prox = [
    "Validar a lista priorizada de 2024 com as equipes pedagógica e psicossocial.",
    "Usar o app no início do próximo ciclo e comparar o previsto com o observado.",
    "Reavaliar o modelo a cada nova rodada da PEDE.",
]
yy = 3.6
for t in prox:
    cx_texto(s, 1.0, yy, 11.2, 0.5, "▪  " + t, 16, RGBColor(0xC5, 0xD3, 0xE8))
    yy += 0.62
cx_texto(s, 0.9, 6.3, 11.5, 0.6,
         f"Aplicação: {APP_URL}  ·  Código, notebook e análises no repositório do projeto.",
         12, DOURADO)

prs.save(OUT)
print("Apresentacao salva em:", OUT, "| slides:", len(prs.slides._sldIdLst))
